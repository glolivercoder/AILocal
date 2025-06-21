#!/usr/bin/env python3
"""
Sistema de Aprimoramento de Conhecimento Constante
Integração com LangChain, TensorFlow e múltiplos formatos
"""

import os
import json
import logging
import asyncio
from typing import Dict, List, Any, Optional, Union
from pathlib import Path
from datetime import datetime
import threading
import queue

# Dependências principais
try:
    import langchain
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain_community.embeddings import OpenAIEmbeddings, HuggingFaceEmbeddings
    from langchain_community.vectorstores import FAISS, Chroma
    from langchain_community.document_loaders import (
        PyPDFLoader, TextLoader, Docx2txtLoader, 
        UnstructuredEPubLoader, UnstructuredPowerPointLoader,
        UnstructuredExcelLoader, UnstructuredWordDocumentLoader
    )
    from langchain.chains import RetrievalQA, ConversationalRetrievalChain
    from langchain.memory import ConversationBufferMemory
    from langchain_community.llms import OpenAI, HuggingFaceHub
    LANGCHAIN_AVAILABLE = True
except ImportError:
    LANGCHAIN_AVAILABLE = False
    print("⚠️ LangChain não disponível. Instale com: pip install langchain")

# TensorFlow para análise de dados
try:
    import tensorflow as tf
    from tensorflow import keras
    from tensorflow.keras import layers
    import numpy as np
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.cluster import KMeans
    TENSORFLOW_AVAILABLE = True
except ImportError:
    TENSORFLOW_AVAILABLE = False
    print("⚠️ TensorFlow não disponível. Instale com: pip install tensorflow scikit-learn")

# Processamento de documentos
try:
    import docx2txt
    import PyPDF2
    import ebooklib
    from ebooklib import epub
    import pandas as pd
    import openpyxl
    from pptx import Presentation
    DOCUMENT_PROCESSING_AVAILABLE = True
except ImportError:
    DOCUMENT_PROCESSING_AVAILABLE = False
    print("⚠️ Bibliotecas de processamento não disponíveis")

# Configurar logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Processador de documentos para múltiplos formatos"""
    
    def __init__(self):
        self.supported_formats = {
            # Microsoft Office
            '.docx': self._process_docx,
            '.doc': self._process_doc,
            '.xlsx': self._process_xlsx,
            '.xls': self._process_xls,
            '.pptx': self._process_pptx,
            '.ppt': self._process_ppt,
            
            # LibreOffice
            '.odt': self._process_odt,
            '.ods': self._process_ods,
            '.odp': self._process_odp,
            
            # E-books
            '.epub': self._process_epub,
            '.mobi': self._process_mobi,
            
            # Outros
            '.pdf': self._process_pdf,
            '.txt': self._process_txt,
            '.md': self._process_markdown,
            '.csv': self._process_csv
        }
    
    def process_document(self, file_path: str) -> Dict[str, Any]:
        """Processa documento e extrai conteúdo"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"Arquivo não encontrado: {file_path}")
        
        file_extension = file_path.suffix.lower()
        
        if file_extension not in self.supported_formats:
            raise ValueError(f"Formato não suportado: {file_extension}")
        
        try:
            processor = self.supported_formats[file_extension]
            content = processor(file_path)
            
            return {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_type": file_extension,
                "content": content,
                "processed_at": datetime.now().isoformat(),
                "status": "success"
            }
            
        except Exception as e:
            logger.error(f"Erro ao processar {file_path}: {e}")
            return {
                "file_path": str(file_path),
                "file_name": file_path.name,
                "file_type": file_extension,
                "content": "",
                "error": str(e),
                "processed_at": datetime.now().isoformat(),
                "status": "error"
            }
    
    def _process_pdf(self, file_path: Path) -> str:
        """Processa arquivo PDF"""
        try:
            if LANGCHAIN_AVAILABLE:
                loader = PyPDFLoader(str(file_path))
                pages = loader.load()
                return "\n".join([page.page_content for page in pages])
            else:
                # Fallback para PyPDF2
                with open(file_path, 'rb') as file:
                    reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                    return text
        except Exception as e:
            logger.error(f"Erro ao processar PDF {file_path}: {e}")
            return ""
    
    def _process_docx(self, file_path: Path) -> str:
        """Processa arquivo DOCX"""
        try:
            if LANGCHAIN_AVAILABLE:
                loader = Docx2txtLoader(str(file_path))
                return loader.load()[0].page_content
            else:
                return docx2txt.process(str(file_path))
        except Exception as e:
            logger.error(f"Erro ao processar DOCX {file_path}: {e}")
            return ""
    
    def _process_doc(self, file_path: Path) -> str:
        """Processa arquivo DOC (legado)"""
        try:
            # Tentar converter DOC para DOCX primeiro
            return docx2txt.process(str(file_path))
        except Exception as e:
            logger.error(f"Erro ao processar DOC {file_path}: {e}")
            return ""
    
    def _process_xlsx(self, file_path: Path) -> str:
        """Processa arquivo Excel"""
        try:
            if LANGCHAIN_AVAILABLE:
                loader = UnstructuredExcelLoader(str(file_path))
                return loader.load()[0].page_content
            else:
                df = pd.read_excel(file_path)
                return df.to_string()
        except Exception as e:
            logger.error(f"Erro ao processar Excel {file_path}: {e}")
            return ""
    
    def _process_xls(self, file_path: Path) -> str:
        """Processa arquivo Excel legado"""
        try:
            df = pd.read_excel(file_path, engine='xlrd')
            return df.to_string()
        except Exception as e:
            logger.error(f"Erro ao processar XLS {file_path}: {e}")
            return ""
    
    def _process_pptx(self, file_path: Path) -> str:
        """Processa arquivo PowerPoint"""
        try:
            if LANGCHAIN_AVAILABLE:
                loader = UnstructuredPowerPointLoader(str(file_path))
                return loader.load()[0].page_content
            else:
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
        except Exception as e:
            logger.error(f"Erro ao processar PowerPoint {file_path}: {e}")
            return ""
    
    def _process_ppt(self, file_path: Path) -> str:
        """Processa arquivo PowerPoint legado"""
        # Implementar conversão ou usar biblioteca específica
        return self._process_pptx(file_path)
    
    def _process_odt(self, file_path: Path) -> str:
        """Processa arquivo ODT (LibreOffice)"""
        try:
            # LibreOffice pode ser convertido para DOCX
            return docx2txt.process(str(file_path))
        except Exception as e:
            logger.error(f"Erro ao processar ODT {file_path}: {e}")
            return ""
    
    def _process_ods(self, file_path: Path) -> str:
        """Processa arquivo ODS (LibreOffice)"""
        try:
            df = pd.read_excel(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Erro ao processar ODS {file_path}: {e}")
            return ""
    
    def _process_odp(self, file_path: Path) -> str:
        """Processa arquivo ODP (LibreOffice)"""
        # Similar ao PowerPoint
        return self._process_pptx(file_path)
    
    def _process_epub(self, file_path: Path) -> str:
        """Processa arquivo EPUB"""
        try:
            if LANGCHAIN_AVAILABLE:
                loader = UnstructuredEPubLoader(str(file_path))
                return loader.load()[0].page_content
            else:
                book = epub.read_epub(file_path)
                text = ""
                for item in book.get_items():
                    if item.get_type() == ebooklib.ITEM_DOCUMENT:
                        text += item.get_content().decode('utf-8') + "\n"
                return text
        except Exception as e:
            logger.error(f"Erro ao processar EPUB {file_path}: {e}")
            return ""
    
    def _process_mobi(self, file_path: Path) -> str:
        """Processa arquivo MOBI (Kindle)"""
        try:
            # Converter MOBI para EPUB primeiro
            # Implementar conversão ou usar biblioteca específica
            return "MOBI processing not implemented yet"
        except Exception as e:
            logger.error(f"Erro ao processar MOBI {file_path}: {e}")
            return ""
    
    def _process_txt(self, file_path: Path) -> str:
        """Processa arquivo de texto"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            logger.error(f"Erro ao processar TXT {file_path}: {e}")
            return ""
    
    def _process_markdown(self, file_path: Path) -> str:
        """Processa arquivo Markdown"""
        return self._process_txt(file_path)
    
    def _process_csv(self, file_path: Path) -> str:
        """Processa arquivo CSV"""
        try:
            df = pd.read_csv(file_path)
            return df.to_string()
        except Exception as e:
            logger.error(f"Erro ao processar CSV {file_path}: {e}")
            return ""

class TensorFlowAnalyzer:
    """Analisador de dados usando TensorFlow"""
    
    def __init__(self):
        self.model = None
        self.vectorizer = None
        self.cluster_model = None
        self.is_trained = False
    
    def create_text_analysis_model(self, vocab_size: int = 10000, max_length: int = 100):
        """Cria modelo de análise de texto"""
        if not TENSORFLOW_AVAILABLE:
            raise ImportError("TensorFlow não disponível")
        
        self.model = keras.Sequential([
            layers.Embedding(vocab_size, 128, input_length=max_length),
            layers.Bidirectional(layers.LSTM(64, return_sequences=True)),
            layers.Bidirectional(layers.LSTM(32)),
            layers.Dense(64, activation='relu'),
            layers.Dropout(0.5),
            layers.Dense(1, activation='sigmoid')
        ])
        
        self.model.compile(
            optimizer='adam',
            loss='binary_crossentropy',
            metrics=['accuracy']
        )
        
        logger.info("Modelo de análise de texto criado")
    
    def extract_features(self, texts: List[str]) -> np.ndarray:
        """Extrai características dos textos"""
        if not TENSORFLOW_AVAILABLE:
            return np.array([])
        
        if self.vectorizer is None:
            self.vectorizer = TfidfVectorizer(
                max_features=1000,
                stop_words='english',
                ngram_range=(1, 2)
            )
        
        features = self.vectorizer.fit_transform(texts)
        return features.toarray()
    
    def cluster_documents(self, texts: List[str], n_clusters: int = 5) -> Dict[str, Any]:
        """Agrupa documentos por similaridade"""
        if not TENSORFLOW_AVAILABLE:
            return {"error": "TensorFlow não disponível"}
        
        try:
            # Extrair características
            features = self.extract_features(texts)
            
            # Aplicar clustering
            self.cluster_model = KMeans(n_clusters=n_clusters, random_state=42)
            clusters = self.cluster_model.fit_predict(features)
            
            # Organizar resultados
            results = {}
            for i, cluster_id in enumerate(clusters):
                if f"cluster_{cluster_id}" not in results:
                    results[f"cluster_{cluster_id}"] = []
                results[f"cluster_{cluster_id}"].append({
                    "text": texts[i][:100] + "...",
                    "index": i
                })
            
            return {
                "clusters": results,
                "n_clusters": n_clusters,
                "total_documents": len(texts),
                "cluster_centers": self.cluster_model.cluster_centers_.tolist()
            }
            
        except Exception as e:
            logger.error(f"Erro no clustering: {e}")
            return {"error": str(e)}
    
    def analyze_sentiment(self, texts: List[str]) -> List[Dict[str, Any]]:
        """Análise de sentimento básica"""
        if not TENSORFLOW_AVAILABLE:
            return [{"error": "TensorFlow não disponível"} for _ in texts]
        
        try:
            # Análise simples baseada em palavras-chave
            positive_words = ['good', 'great', 'excellent', 'amazing', 'wonderful', 'perfect']
            negative_words = ['bad', 'terrible', 'awful', 'horrible', 'worst', 'disappointing']
            
            results = []
            for text in texts:
                text_lower = text.lower()
                positive_count = sum(1 for word in positive_words if word in text_lower)
                negative_count = sum(1 for word in negative_words if word in text_lower)
                
                if positive_count > negative_count:
                    sentiment = "positive"
                    score = positive_count / (positive_count + negative_count + 1)
                elif negative_count > positive_count:
                    sentiment = "negative"
                    score = negative_count / (positive_count + negative_count + 1)
                else:
                    sentiment = "neutral"
                    score = 0.5
                
                results.append({
                    "text": text[:100] + "...",
                    "sentiment": sentiment,
                    "score": score,
                    "positive_words": positive_count,
                    "negative_words": negative_count
                })
            
            return results
            
        except Exception as e:
            logger.error(f"Erro na análise de sentimento: {e}")
            return [{"error": str(e)} for _ in texts]

class LangChainEnhancer:
    """Sistema de aprimoramento usando LangChain"""
    
    def __init__(self, openai_api_key: str = None, model_name: str = "gpt-3.5-turbo"):
        self.openai_api_key = openai_api_key
        self.model_name = model_name
        self.embeddings = None
        self.vectorstore = None
        self.qa_chain = None
        self.memory = None
        
        if not LANGCHAIN_AVAILABLE:
            logger.warning("LangChain não disponível")
            return
        
        self._initialize_components()
    
    def _initialize_components(self):
        """Inicializa componentes do LangChain"""
        try:
            # Configurar embeddings
            if self.openai_api_key:
                self.embeddings = OpenAIEmbeddings(openai_api_key=self.openai_api_key)
            else:
                # Usar embeddings locais
                self.embeddings = HuggingFaceEmbeddings(
                    model_name="sentence-transformers/all-MiniLM-L6-v2"
                )
            
            # Configurar memória
            self.memory = ConversationBufferMemory(
                memory_key="chat_history",
                return_messages=True
            )
            
            logger.info("Componentes LangChain inicializados")
            
        except Exception as e:
            logger.error(f"Erro ao inicializar LangChain: {e}")
    
    def create_vectorstore(self, documents: List[str], persist_directory: str = "vectorstore"):
        """Cria vectorstore com documentos"""
        if not LANGCHAIN_AVAILABLE or not self.embeddings:
            return False
        
        try:
            # Dividir documentos
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200
            )
            
            texts = []
            for doc in documents:
                if isinstance(doc, str):
                    texts.extend(text_splitter.split_text(doc))
                else:
                    texts.extend(text_splitter.split_text(doc.page_content))
            
            # Criar vectorstore
            self.vectorstore = FAISS.from_texts(texts, self.embeddings)
            
            # Salvar
            self.vectorstore.save_local(persist_directory)
            
            logger.info(f"Vectorstore criado com {len(texts)} chunks")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao criar vectorstore: {e}")
            return False
    
    def load_vectorstore(self, persist_directory: str = "vectorstore"):
        """Carrega vectorstore existente"""
        if not LANGCHAIN_AVAILABLE or not self.embeddings:
            return False
        
        try:
            self.vectorstore = FAISS.load_local(persist_directory, self.embeddings)
            logger.info("Vectorstore carregado")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao carregar vectorstore: {e}")
            return False
    
    def query_knowledge_base(self, question: str, k: int = 4) -> Dict[str, Any]:
        """Consulta a base de conhecimento"""
        if not LANGCHAIN_AVAILABLE or not self.vectorstore:
            return {"error": "LangChain ou vectorstore não disponível"}
        
        try:
            # Buscar documentos relevantes
            docs = self.vectorstore.similarity_search(question, k=k)
            
            # Criar chain de QA se não existir
            if not self.qa_chain:
                from langchain.chat_models import ChatOpenAI
                
                llm = ChatOpenAI(
                    model_name=self.model_name,
                    openai_api_key=self.openai_api_key
                )
                
                self.qa_chain = ConversationalRetrievalChain.from_llm(
                    llm=llm,
                    retriever=self.vectorstore.as_retriever(),
                    memory=self.memory
                )
            
            # Executar query
            result = self.qa_chain({"question": question})
            
            return {
                "question": question,
                "answer": result["answer"],
                "source_documents": [doc.page_content[:200] + "..." for doc in docs],
                "chat_history": result.get("chat_history", [])
            }
            
        except Exception as e:
            logger.error(f"Erro na consulta: {e}")
            return {"error": str(e)}

class KnowledgeEnhancementSystem:
    """Sistema principal de aprimoramento de conhecimento"""
    
    def __init__(self, config_file: str = "config/knowledge_config.json"):
        self.config_file = config_file
        self.config = self.load_config()
        
        # Inicializar componentes
        self.document_processor = DocumentProcessor()
        self.tensorflow_analyzer = TensorFlowAnalyzer()
        self.langchain_enhancer = LangChainEnhancer(
            openai_api_key=self.config.get("openai_api_key")
        )
        
        # Estado do sistema
        self.processed_documents = []
        self.knowledge_base = []
        self.analysis_results = {}
        
        # Threading para processamento assíncrono
        self.processing_queue = queue.Queue()
        self.processing_thread = None
        self.is_processing = False
        
        logger.info("Sistema de Aprimoramento de Conhecimento inicializado")
    
    def load_config(self) -> Dict[str, Any]:
        """Carrega configuração do sistema"""
        try:
            if Path(self.config_file).exists():
                with open(self.config_file, 'r') as f:
                    return json.load(f)
            else:
                # Configuração padrão
                default_config = {
                    "openai_api_key": "",
                    "model_name": "gpt-3.5-turbo",
                    "vectorstore_path": "vectorstore",
                    "max_documents": 1000,
                    "chunk_size": 1000,
                    "chunk_overlap": 200,
                    "enable_tensorflow": TENSORFLOW_AVAILABLE,
                    "enable_langchain": LANGCHAIN_AVAILABLE
                }
                
                # Salvar configuração padrão
                Path(self.config_file).parent.mkdir(parents=True, exist_ok=True)
                with open(self.config_file, 'w') as f:
                    json.dump(default_config, f, indent=2)
                
                return default_config
                
        except Exception as e:
            logger.error(f"Erro ao carregar configuração: {e}")
            return {}
    
    def start_processing(self):
        """Inicia processamento em background"""
        if not self.is_processing:
            self.is_processing = True
            self.processing_thread = threading.Thread(target=self._processing_worker)
            self.processing_thread.daemon = True
            self.processing_thread.start()
            logger.info("Processamento em background iniciado")
    
    def stop_processing(self):
        """Para processamento em background"""
        self.is_processing = False
        if self.processing_thread:
            self.processing_thread.join()
        logger.info("Processamento em background parado")
    
    def _processing_worker(self):
        """Worker para processamento em background"""
        while self.is_processing:
            try:
                # Processar itens da fila
                if not self.processing_queue.empty():
                    task = self.processing_queue.get(timeout=1)
                    self._process_task(task)
                else:
                    time.sleep(1)
            except queue.Empty:
                continue
            except Exception as e:
                logger.error(f"Erro no worker: {e}")
    
    def _process_task(self, task: Dict[str, Any]):
        """Processa uma tarefa"""
        task_type = task.get("type")
        
        if task_type == "document":
            self._process_document_task(task)
        elif task_type == "analysis":
            self._process_analysis_task(task)
        elif task_type == "enhancement":
            self._process_enhancement_task(task)
    
    def _process_document_task(self, task: Dict[str, Any]):
        """Processa tarefa de documento"""
        file_path = task.get("file_path")
        
        # Processar documento
        result = self.document_processor.process_document(file_path)
        
        if result["status"] == "success":
            self.processed_documents.append(result)
            self.knowledge_base.append(result["content"])
            
            # Adicionar análise à fila
            self.processing_queue.put({
                "type": "analysis",
                "content": result["content"],
                "file_info": result
            })
    
    def _process_analysis_task(self, task: Dict[str, Any]):
        """Processa tarefa de análise"""
        content = task.get("content")
        
        if TENSORFLOW_AVAILABLE:
            # Análise com TensorFlow
            analysis = self.tensorflow_analyzer.analyze_sentiment([content])
            self.analysis_results[task.get("file_info", {}).get("file_name", "unknown")] = analysis[0]
    
    def _process_enhancement_task(self, task: Dict[str, Any]):
        """Processa tarefa de aprimoramento"""
        question = task.get("question")
        
        if LANGCHAIN_AVAILABLE:
            # Consulta com LangChain
            result = self.langchain_enhancer.query_knowledge_base(question)
            task["result"] = result
    
    def add_document(self, file_path: str) -> Dict[str, Any]:
        """Adiciona documento para processamento"""
        if not Path(file_path).exists():
            return {"error": "Arquivo não encontrado"}
        
        # Adicionar à fila de processamento
        self.processing_queue.put({
            "type": "document",
            "file_path": file_path
        })
        
        return {
            "status": "queued",
            "file_path": file_path,
            "message": "Documento adicionado à fila de processamento"
        }
    
    def query_knowledge(self, question: str) -> Dict[str, Any]:
        """Consulta a base de conhecimento"""
        if LANGCHAIN_AVAILABLE:
            return self.langchain_enhancer.query_knowledge_base(question)
        else:
            return {"error": "LangChain não disponível"}
    
    def analyze_documents(self, n_clusters: int = 5) -> Dict[str, Any]:
        """Analisa documentos com TensorFlow"""
        if not TENSORFLOW_AVAILABLE:
            return {"error": "TensorFlow não disponível"}
        
        if not self.knowledge_base:
            return {"error": "Nenhum documento processado"}
        
        return self.tensorflow_analyzer.cluster_documents(self.knowledge_base, n_clusters)
    
    def get_system_status(self) -> Dict[str, Any]:
        """Retorna status do sistema"""
        return {
            "processed_documents": len(self.processed_documents),
            "knowledge_base_size": len(self.knowledge_base),
            "analysis_results": len(self.analysis_results),
            "queue_size": self.processing_queue.qsize(),
            "is_processing": self.is_processing,
            "tensorflow_available": TENSORFLOW_AVAILABLE,
            "langchain_available": LANGCHAIN_AVAILABLE,
            "document_processing_available": DOCUMENT_PROCESSING_AVAILABLE
        }
    
    def save_knowledge_base(self, file_path: str = "knowledge_base.json"):
        """Salva base de conhecimento"""
        try:
            data = {
                "processed_documents": self.processed_documents,
                "knowledge_base": self.knowledge_base,
                "analysis_results": self.analysis_results,
                "saved_at": datetime.now().isoformat()
            }
            
            with open(file_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
            
            logger.info(f"Base de conhecimento salva em {file_path}")
            return True
            
        except Exception as e:
            logger.error(f"Erro ao salvar base de conhecimento: {e}")
            return False
    
    def load_knowledge_base(self, file_path: str = "knowledge_base.json"):
        """Carrega base de conhecimento"""
        try:
            if Path(file_path).exists():
                with open(file_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                
                self.processed_documents = data.get("processed_documents", [])
                self.knowledge_base = data.get("knowledge_base", [])
                self.analysis_results = data.get("analysis_results", {})
                
                logger.info(f"Base de conhecimento carregada de {file_path}")
                return True
            else:
                logger.warning(f"Arquivo de base de conhecimento não encontrado: {file_path}")
                return False
                
        except Exception as e:
            logger.error(f"Erro ao carregar base de conhecimento: {e}")
            return False

def main():
    """Função principal para teste do sistema"""
    print("🧠 Iniciando Sistema de Aprimoramento de Conhecimento...")
    
    # Criar sistema
    system = KnowledgeEnhancementSystem()
    
    # Iniciar processamento
    system.start_processing()
    
    # Testar processamento de documentos
    test_files = [
        "test_document.txt",
        "test_document.pdf",
        "test_document.docx"
    ]
    
    print("\n📄 Testando processamento de documentos:")
    for file_path in test_files:
        if Path(file_path).exists():
            result = system.add_document(file_path)
            print(f"  {file_path}: {result}")
    
    # Aguardar processamento
    time.sleep(5)
    
    # Status do sistema
    print("\n📊 Status do Sistema:")
    status = system.get_system_status()
    for key, value in status.items():
        print(f"  {key}: {value}")
    
    # Testar consulta
    if LANGCHAIN_AVAILABLE:
        print("\n🔍 Testando consulta:")
        result = system.query_knowledge("O que é inteligência artificial?")
        print(f"Resposta: {result.get('answer', 'Erro na consulta')}")
    
    # Testar análise
    if TENSORFLOW_AVAILABLE:
        print("\n📈 Testando análise:")
        result = system.analyze_documents(n_clusters=3)
        print(f"Clusters: {len(result.get('clusters', {}))}")
    
    # Parar processamento
    system.stop_processing()
    print("\n✅ Teste concluído!")

if __name__ == "__main__":
    main() 